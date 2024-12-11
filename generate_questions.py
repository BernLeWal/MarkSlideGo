#!/usr/bin/env python3
""" Generate the Moodle SC/MC-questions for a course from a YAML file. """
import os
import sys
import logging
from enum import Enum
from openai import AzureOpenAI
#from azure.identity import DefaultAzureCredential, get_bearer_token_provider
from dotenv import load_dotenv
from pptx import Presentation
import fitz # PyMuPDF


load_dotenv()  # take environment variables from .env.
# Setup logging
LOGGING_LEVEL = os.environ.get('LOGGING_LEVEL', 'INFO')  # default to INFO if no env var set
numeric_level = getattr(logging, LOGGING_LEVEL.upper(), None)
if not isinstance(numeric_level, int):
    raise ValueError(f'Invalid log level: {LOGGING_LEVEL}')
logging.basicConfig(level=numeric_level)
logger = logging.getLogger(__name__)


endpoint = os.getenv("OPENAI_ENDPOINT_URL")
apikey = os.getenv("OPENAI_API_KEY")
deployment = os.getenv("OPENAI_DEPLOYMENT_NAME")



class QuestionFormat(Enum):
    """ Enumeration of question formats with samples. """
    QFORMAT = """
Use the following output structure (this first line contains the table column headers, every further line represents a question; Question1 and Question2 are samples not to be contained in your output):
questionname,questiontext,A,B,C,D,Answer 1,Answer 2,answernumbering,correctfeedback,partiallycorrectfeedback,incorrectfeedback,defaultmark
Question1,The dmesg command,Shows user login logoff attempts,Shows the syslog file for info messages,kernel log messages,Shows the daemon log messages,C,,123,Your answer is correct.,Your answer is partially correct.,Your answer is incorrect.,1
Question2,The command “mknod myfifo b 4 16”,Will create a block device if user is root,Will create a block device for all users,Will create a FIFO if user is not root,"None ,of the mentioned",A,B,ABCD,Your answer is correct.,Your answer is partially correct.,Your answer is incorrect.,1
    """
    # AIKEN format description see https://docs.moodle.org/403/en/Aiken_format
    AIKEN = """
What is the correct answer to this question?
A. This is not the correct answer
B. This answer is wrong
C. Also here a wrong answer
D. This answer is correct!
ANSWER: D    
    """
    # GIFT format description see https://docs.moodle.org/403/en/GIFT_format, https://docs.moodle.org/404/en/GIFT_format#Multiple_Answers
    GIFT = """
//A Comment for a question1
:: What is the correct answer to this question?{
=A correct answer
~Wrong answer1
#A response to wrong answer1
~Wrong answer2
#A response to wrong answer2
~Wrong answer3
#A response to wrong answer3
~Wrong answer4
#A response to wrong answer4
}
    """


def load_contents(source_file:str, export_content:bool=True)->str:
    """ Load contents of a file into a string.  """
    logger.info("Loading contents of %s ...", source_file)
    content = ""
    source_filename, source_ext = os.path.splitext(source_file)
    if source_ext == '.pptx':
        # see https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
        presentation = Presentation(source_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
    elif source_ext == '.pdf':
        # see https://pymupdf.readthedocs.io/en/latest/tutorial.html
        pdf_document = fitz.open(source_file)
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            content += page.get_text()
    else:
        # treat as text-file
        with open(source_file, 'r', encoding="utf-8") as content_file:
            content = content_file.read()
    if export_content:
        with open(source_filename + '.txt', 'w', encoding="utf-8") as output_file:
            output_file.write(content)
    return content


def generate_questions(course_title:str, slides_title:str, source_file:str, num:int=10,
    target_file:str=None, question_format:QuestionFormat=QuestionFormat.AIKEN)->str:
    """ Generate the Moodle SC/MC-questions for a course from a YAML file. """
    logger.info("Generating questions for %s ...", source_file)
    #token_provider = get_bearer_token_provider(
    #    DefaultAzureCredential(),
    #    "https://cognitiveservices.azure.com/.default")

    client = AzureOpenAI(
        api_key=apikey,
        api_version="2024-05-01-preview",
        azure_endpoint=endpoint,
    #    azure_ad_token_provider=token_provider,
    )

    # Load contents of markdown_file into a string
    content = load_contents(source_file)

    messages= [
    {
        "role": "system",
        "content": f"""You are an tutor for a course on informatics dealing with {course_title}.
The current topic is {slides_title}.
You will create questions about the topic and output as file, which is later used by the Moodle Question import.:
{question_format.value}
"""},{
        "role": "user",
        "content": f"Generate {num} multiple-choice questions with one single correct answere covering only the specified contents, provided in Markdown format, as follows: {content}"
    }]
    #print(messages)

    completion = client.chat.completions.create(
        model=deployment,
        messages=messages,
        max_tokens=1000,    # ca. 100 tokens needed per question
        temperature=0.7,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None,
        stream=False
    )
    questions = completion.choices[0].message.content

    if target_file is not None:
        logger.info("Writing questions to %s", target_file)
        with open(target_file, 'w', encoding="utf-8") as output_file:
            output_file.write(questions)

    return questions



if __name__ == "__main__":
    # Check if two arguments were provided
    if len(sys.argv) < 5:
        print(f"Usage: {sys.argv[0]} <course-title> <slides-title> <markdown-file> <num-questions> [target-file]")
        print()
        print("Example:")
        print(f"{sys.argv[0]} \"Software Engineering 1 - Labor\" \"Introduction to Programming in C#\" \"catalogs/tw-kb/programming/csharp/csharp_intro.md\" 10")
        sys.exit(0)

    # Path to the YAML file
    coursetitle = sys.argv[1]
    slidestitle = sys.argv[2]
    markdownfile = sys.argv[3]
    numquestions = int(sys.argv[4])
    if len(sys.argv) > 5:
        targetfile = sys.argv[5]
        generate_questions(coursetitle, slidestitle, markdownfile, numquestions, targetfile)
        print(f"Questions written to {targetfile}")
    else:
        result = generate_questions(coursetitle, slidestitle, markdownfile, num=numquestions)
        print(result)
